import re
import os
import sys
import tempfile
import argparse
from PIL import Image
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Presentation Theme Colors ---
THEMES = {
    'dark': {
        'bg': RGBColor(15, 23, 42),       # Slate 900
        'text': RGBColor(248, 250, 252),  # Slate 100
        'text_muted': RGBColor(148, 163, 184), # Slate 400
        'primary': RGBColor(99, 102, 241), # Indigo 500
        'block_bg': RGBColor(30, 41, 59), # Slate 800
        'block_border': RGBColor(71, 85, 105), # Slate 600
        'alert_bg': RGBColor(69, 10, 10), # Soft dark red
        'alert_border': RGBColor(185, 28, 28), # Red 700
        'eq_color': 'white'
    },
    'light': {
        'bg': RGBColor(248, 250, 252),    # Slate 50
        'text': RGBColor(15, 23, 42),      # Slate 900
        'text_muted': RGBColor(71, 85, 105), # Slate 600
        'primary': RGBColor(37, 99, 235),  # Blue 600
        'block_bg': RGBColor(241, 245, 249), # Slate 100
        'block_border': RGBColor(203, 213, 225), # Slate 300
        'alert_bg': RGBColor(254, 242, 242), # Red 50
        'alert_border': RGBColor(239, 68, 68), # Red 500
        'eq_color': 'black'
    }
}

# --- Helper function to find matching curly brace content ---
def find_curly_content(text, start_idx):
    if start_idx >= len(text) or text[start_idx] != '{':
        return None, start_idx
    count = 0
    for i in range(start_idx, len(text)):
        if text[i] == '{' and (i == 0 or text[i-1] != '\\'):
            count += 1
        elif text[i] == '}' and (i == 0 or text[i-1] != '\\'):
            count -= 1
            if count == 0:
                return text[start_idx+1:i], i+1
    return None, start_idx

# --- Clean LaTeX formatting command macros ---
def clean_latex_formatting(text):
    if not text:
        return ""
    text = text.replace('\n', ' ').replace('\r', ' ')
    text = text.replace(r'\\', ' ')
    while True:
        match = re.search(r'\\[a-zA-Z]+\*?(?:\[[^\]]*\])?{', text)
        if not match:
            break
        start_idx = match.start()
        open_brace_idx = text.find('{', start_idx)
        content, end_idx = find_curly_content(text, open_brace_idx)
        if content is not None:
            macro_call = text[start_idx:open_brace_idx]
            if any(cmd in macro_call for cmd in ['\\textbf', '\\textit', '\\emph', '\\textcolor', '\\text']):
                replacement = content
            else:
                if '\\titlepage' in macro_call or '\\tableofcontents' in macro_call:
                    replacement = ""
                else:
                    replacement = content
            text = text[:start_idx] + replacement + text[end_idx:]
        else:
            break
            
    text = re.sub(r'\\([%&_{}])', r'\1', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

# --- Tokenize inline text for rich runs (bold, italic, inline math) ---
def tokenize_inline_text(text):
    tokens = []
    pos = 0
    while pos < len(text):
        match = re.search(r'(\$[^\$]+\$|\\[a-zA-Z]+(?:\[[^\]]*\])?{)', text[pos:])
        if not match:
            remaining = text[pos:]
            if remaining:
                tokens.append((remaining, {}))
            break
            
        start_idx = pos + match.start()
        before = text[pos:start_idx]
        if before:
            tokens.append((before, {}))
            
        match_str = match.group(0)
        
        if match_str.startswith('$'):
            math_content = match_str[1:-1]
            tokens.append((math_content, {'math': True, 'italic': True}))
            pos = start_idx + len(match_str)
        else:
            cmd_match = re.match(r'\\([a-zA-Z]+)', match_str)
            cmd = cmd_match.group(1)
            open_brace_idx = start_idx + len(match_str) - 1
            content, end_idx = find_curly_content(text, open_brace_idx)
            
            if content is not None:
                style = {}
                if cmd in ['textbf', 'bold']:
                    style['bold'] = True
                elif cmd in ['textit', 'emph', 'italic']:
                    style['italic'] = True
                elif cmd == 'textcolor':
                    color_match = re.search(r'\\textcolor{([a-zA-Z]+)}', text[start_idx:open_brace_idx+1])
                    if color_match:
                        style['color'] = color_match.group(1).lower()
                
                sub_tokens = tokenize_inline_text(content)
                for text_seg, sub_style in sub_tokens:
                    merged_style = {**style, **sub_style}
                    if 'color' in style and 'color' not in sub_style:
                        merged_style['color'] = style['color']
                    tokens.append((text_seg, merged_style))
                pos = end_idx
            else:
                tokens.append((match_str, {}))
                pos = start_idx + len(match_str)
                
    return tokens

# --- Add formatting runs to presentation paragraphs ---
def add_runs_to_paragraph(paragraph, tokens, theme, font_size=18):
    paragraph.font.name = 'Calibri'
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = theme['text']
    
    for text_seg, style in tokens:
        if not text_seg.strip() and not text_seg == " ":
            continue
        run = paragraph.add_run()
        run.text = text_seg
        run.font.name = 'Calibri'
        run.font.size = Pt(font_size)
        
        if style.get('bold'):
            run.font.bold = True
        if style.get('italic'):
            run.font.italic = True
            
        if style.get('color'):
            c = style.get('color')
            if c == 'red':
                run.font.color.rgb = RGBColor(239, 68, 68)
            elif c == 'blue':
                run.font.color.rgb = RGBColor(59, 130, 246)
            elif c == 'green':
                run.font.color.rgb = RGBColor(34, 197, 94)
            else:
                run.font.color.rgb = theme['primary']
        elif style.get('math'):
            run.font.italic = True
            run.font.color.rgb = theme['primary']
        else:
            run.font.color.rgb = theme['text']

# --- Math equation rendering to transparent PNG ---
def render_latex_to_png(latex_str, text_color='white', dpi=300):
    math_str = latex_str.strip()
    if math_str.startswith('$$') and math_str.endswith('$$'):
        math_str = math_str[2:-2].strip()
    elif math_str.startswith('$') and math_str.endswith('$'):
        math_str = math_str[1:-1].strip()
        
    full_math_str = f"${math_str}$"
    fig = plt.figure(figsize=(0.1, 0.1), dpi=dpi)
    
    # Render with Computer Modern style using mathtext
    plt.text(0.5, 0.5, full_math_str, size=15, color=text_color,
             horizontalalignment='center', verticalalignment='center',
             usetex=False, math_fontfamily='cm')
    plt.axis('off')
    
    temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    temp_path = temp_file.name
    temp_file.close()
    
    try:
        plt.savefig(temp_path, bbox_inches='tight', pad_inches=0.08, transparent=True, dpi=dpi)
        plt.close(fig)
        return temp_path
    except Exception as e:
        plt.close(fig)
        if os.path.exists(temp_path):
            os.remove(temp_path)
        raise e

# --- Recursive parser for LaTeX environments ---
def find_matching_environment(text, env_name, start_idx):
    begin_str = f'\\begin{{{env_name}}}'
    end_str = f'\\end{{{env_name}}}'
    count = 1
    pos = start_idx + len(begin_str)
    
    while pos < len(text):
        next_begin = text.find(begin_str, pos)
        next_end = text.find(end_str, pos)
        
        if next_end == -1:
            return -1
            
        if next_begin != -1 and next_begin < next_end:
            count += 1
            pos = next_begin + len(begin_str)
        else:
            count -= 1
            if count == 0:
                return next_end
            pos = next_end + len(end_str)
            
    return -1

def parse_elements(text):
    text = text.strip()
    elements = []
    pos = 0
    
    while pos < len(text):
        match = re.search(r'(\\begin{[a-zA-Z]+}|\\\[|$$|\$\$)', text[pos:])
        if not match:
            remaining = text[pos:].strip()
            if remaining:
                elements.append({'type': 'paragraph', 'text': remaining})
            break
            
        before = text[pos:pos+match.start()].strip()
        if before:
            elements.append({'type': 'paragraph', 'text': before})
            
        match_str = match.group(0)
        start_idx = pos + match.start()
        
        if match_str in ['$$', '$$', r'\[']:
            end_marker = r'\]' if match_str == r'\[' else '$$'
            eq_end = text.find(end_marker, start_idx + len(match_str))
            if eq_end == -1:
                elements.append({'type': 'paragraph', 'text': text[start_idx:]})
                break
            eq_content = text[start_idx+len(match_str):eq_end].strip()
            elements.append({'type': 'equation', 'latex': eq_content})
            pos = eq_end + len(end_marker)
        elif match_str.startswith(r'\begin'):
            env_name_match = re.match(r'\\begin{([a-zA-Z]+)}', match_str)
            env_name = env_name_match.group(1)
            
            end_tag = f'\\end{{{env_name}}}'
            end_match = find_matching_environment(text, env_name, start_idx)
            if end_match == -1:
                elements.append({'type': 'paragraph', 'text': text[start_idx:]})
                break
                
            env_body = text[start_idx + len(match_str):end_match].strip()
            
            if env_name in ['block', 'alertblock']:
                title_start = start_idx + len(f'\\begin{{{env_name}}}')
                while title_start < len(text) and text[title_start].isspace():
                    title_start += 1
                block_title = ""
                body_start = title_start
                if title_start < len(text) and text[title_start] == '{':
                    title_content, body_start = find_curly_content(text, title_start)
                    if title_content is not None:
                        block_title = title_content
                
                block_body = text[body_start:end_match].strip()
                elements.append({
                    'type': 'block',
                    'block_type': 'standard' if env_name == 'block' else 'alert',
                    'title': clean_latex_formatting(block_title),
                    'children': parse_elements(block_body)
                })
                pos = end_match + len(end_tag)
                
            elif env_name == 'columns':
                elements.append({
                    'type': 'columns',
                    'columns': parse_columns(env_body)
                })
                pos = end_match + len(end_tag)
                
            elif env_name == 'itemize':
                elements.append({
                    'type': 'itemize',
                    'items': parse_itemize(env_body)
                })
                pos = end_match + len(end_tag)
                
            elif env_name == 'table':
                elements.extend(parse_elements(env_body))
                pos = end_match + len(end_tag)
                
            elif env_name == 'tabular':
                elements.append({
                    'type': 'table',
                    'data': parse_tabular(env_body)
                })
                pos = end_match + len(end_tag)
                
            else:
                elements.extend(parse_elements(env_body))
                pos = end_match + len(end_tag)
                
    return elements

def parse_itemize(body):
    items = []
    pos = 0
    item_positions = []
    nesting_level = 0
    
    while pos < len(body):
        if body.startswith(r'\begin{itemize}', pos):
            nesting_level += 1
            pos += len(r'\begin{itemize}')
        elif body.startswith(r'\end{itemize}', pos):
            nesting_level -= 1
            pos += len(r'\end{itemize}')
        elif body.startswith(r'\item', pos) and nesting_level == 0:
            item_positions.append(pos)
            pos += len(r'\item')
        else:
            pos += 1
            
    for i in range(len(item_positions)):
        start = item_positions[i] + len(r'\item')
        end = item_positions[i+1] if i + 1 < len(item_positions) else len(body)
        item_text = body[start:end].strip()
        if item_text.startswith('['):
            close_bracket = item_text.find(']')
            if close_bracket != -1 and close_bracket < 10:
                item_text = item_text[close_bracket+1:].strip()
        items.append(parse_elements(item_text))
        
    return items

def parse_columns(body):
    columns = []
    pos = 0
    while True:
        col_start = body.find(r'\begin{column}', pos)
        if col_start == -1:
            break
        width_brace_start = col_start + len(r'\begin{column}')
        while width_brace_start < len(body) and body[width_brace_start].isspace():
            width_brace_start += 1
        width = 0.5
        body_start = width_brace_start
        if width_brace_start < len(body) and body[width_brace_start] == '{':
            width_content, body_start = find_curly_content(body, width_brace_start)
            if width_content is not None:
                match = re.match(r'([0-9.]+)', width_content)
                if match:
                    width = float(match.group(1))
                    
        col_end = find_matching_environment(body, 'column', col_start)
        if col_end == -1:
            break
            
        col_body = body[body_start:col_end].strip()
        columns.append({
            'width': width,
            'children': parse_elements(col_body)
        })
        pos = col_end + len(r'\end{column}')
        
    return columns

def parse_tabular(body):
    rows = []
    raw_lines = body.split(r'\\')
    for raw_line in raw_lines:
        raw_line = raw_line.strip()
        cleaned_line = re.sub(r'\\(hline|toprule|midrule|bottomrule)', '', raw_line).strip()
        if not cleaned_line:
            continue
        cols = [clean_latex_formatting(c.strip()) for c in cleaned_line.split('&')]
        rows.append(cols)
    return rows

# --- Height Estimator for slide elements ---
def estimate_element_height(element, width_inches, font_size=18):
    line_height = font_size * 1.45 / 72.0
    char_width = 0.075 * (font_size / 18.0)
    chars_per_line = max(10, int(width_inches / char_width))
    
    if element['type'] == 'paragraph':
        text = element['text']
        lines = max(1, len(text) // chars_per_line + 1)
        return lines * line_height + 0.15
        
    elif element['type'] == 'equation':
        return 0.7 + 0.15
        
    elif element['type'] == 'itemize':
        h = 0
        for item in element['items']:
            item_h = 0
            for child in item:
                if child['type'] == 'paragraph':
                    lines = max(1, len(child['text']) // (chars_per_line - 4) + 1)
                    item_h += lines * line_height
                else:
                    item_h += estimate_element_height(child, width_inches - 0.4, font_size)
            h += max(line_height, item_h) + 0.08
        return h + 0.15
        
    elif element['type'] == 'block':
        children_h = 0
        for child in element['children']:
            children_h += estimate_element_height(child, width_inches - 0.4, font_size)
        return 0.4 + children_h + 0.2
        
    elif element['type'] == 'columns':
        col_heights = []
        col_width = (width_inches - 0.4 * (len(element['columns']) - 1)) / len(element['columns'])
        for col in element['columns']:
            col_h = 0
            for child in col['children']:
                col_h += estimate_element_height(child, col_width, font_size)
            col_heights.append(col_h)
        return max(col_heights) if col_heights else 0
        
    elif element['type'] == 'table':
        return len(element['data']) * 0.45 + 0.2
        
    return 0

# --- Render elements vertically flowing onto PPTX Slide ---
def render_elements_to_slide(slide, elements, theme, theme_name, left, width, start_y):
    current_y = start_y
    i = 0
    while i < len(elements):
        elem = elements[i]
        
        if elem['type'] in ['paragraph', 'itemize']:
            text_elems = []
            while i < len(elements) and elements[i]['type'] in ['paragraph', 'itemize']:
                text_elems.append(elements[i])
                i += 1
            
            text_height = 0
            for te in text_elems:
                text_height += estimate_element_height(te, width.inches)
                
            txBox = slide.shapes.add_textbox(left, current_y, width, Inches(text_height))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            
            first = True
            for te in text_elems:
                if te['type'] == 'paragraph':
                    tokens = tokenize_inline_text(te['text'])
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    add_runs_to_paragraph(p, tokens, theme)
                    p.space_after = Pt(8)
                elif te['type'] == 'itemize':
                    for item in te['items']:
                        item_text = ""
                        for child in item:
                            if child['type'] == 'paragraph':
                                item_text += child['text']
                        
                        tokens = tokenize_inline_text(item_text)
                        if first:
                            p = tf.paragraphs[0]
                            first = False
                        else:
                            p = tf.add_paragraph()
                        p.level = 0
                        add_runs_to_paragraph(p, tokens, theme)
                        p.space_after = Pt(6)
            
            current_y += Inches(text_height)
            continue
            
        elif elem['type'] == 'block':
            block_h = Inches(estimate_element_height(elem, width.inches))
            
            block_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, current_y, width, block_h
            )
            block_shape.fill.solid()
            block_shape.fill.fore_color.rgb = theme['block_bg'] if elem['block_type'] == 'standard' else theme['alert_bg']
            block_shape.line.color.rgb = theme['block_border'] if elem['block_type'] == 'standard' else theme['alert_border']
            block_shape.line.width = Pt(1.5)
            
            tb_left = left + Inches(0.15)
            tb_top = current_y + Inches(0.1)
            tb_width = width - Inches(0.3)
            tb_height = block_h - Inches(0.2)
            
            txBox = slide.shapes.add_textbox(tb_left, tb_top, tb_width, tb_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            
            p_title = tf.paragraphs[0]
            p_title.text = elem['title']
            p_title.font.name = 'Calibri'
            p_title.font.size = Pt(20)
            p_title.font.bold = True
            if elem['block_type'] == 'standard':
                p_title.font.color.rgb = theme['primary']
            else:
                p_title.font.color.rgb = theme['alert_border'] if theme_name == 'light' else RGBColor(251, 113, 133)
            p_title.space_after = Pt(8)
            
            for child in elem['children']:
                if child['type'] == 'paragraph':
                    tokens = tokenize_inline_text(child['text'])
                    p = tf.add_paragraph()
                    add_runs_to_paragraph(p, tokens, theme)
                    p.space_after = Pt(6)
                elif child['type'] == 'itemize':
                    for item in child['items']:
                        item_text = ""
                        for grandchild in item:
                            if grandchild['type'] == 'paragraph':
                                item_text += grandchild['text']
                        
                        tokens = tokenize_inline_text(item_text)
                        p = tf.add_paragraph()
                        p.level = 0
                        add_runs_to_paragraph(p, tokens, theme)
                        p.space_after = Pt(4)
                        
            current_y += block_h + Inches(0.15)
            
        elif elem['type'] == 'equation':
            eq_h = Inches(0.7)
            try:
                img_path = render_latex_to_png(elem['latex'], text_color=theme['eq_color'])
                if img_path and os.path.exists(img_path):
                    img = Image.open(img_path)
                    img_w_px, img_h_px = img.size
                    display_h = Inches(0.4)
                    display_w = Inches((img_w_px / img_h_px) * 0.4)
                    if display_w > width:
                        display_w = width
                        display_h = Inches((img_h_px / img_w_px) * width.inches)
                        
                    img_left = left + (width - display_w) / 2
                    img_top = current_y + (eq_h - display_h) / 2
                    
                    slide.shapes.add_picture(img_path, img_left, img_top, width=display_w, height=display_h)
                    try:
                        os.remove(img_path)
                    except:
                        pass
                else:
                    raise Exception("Failed to render equation image")
            except Exception as e:
                print(f"Fallback to text for equation due to: {e}")
                txBox = slide.shapes.add_textbox(left, current_y, width, eq_h)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = elem['latex']
                p.alignment = PP_ALIGN.CENTER
                p.font.name = 'Courier New'
                p.font.size = Pt(16)
                p.font.italic = True
                p.font.color.rgb = theme['primary']
                
            current_y += eq_h + Inches(0.15)
            
        elif elem['type'] == 'columns':
            col_h = Inches(estimate_element_height(elem, width.inches))
            cols = elem['columns']
            N = len(cols)
            if N > 0:
                gap = Inches(0.4)
                col_width = (width - gap * (N - 1)) / N
                
                for c_idx, col in enumerate(cols):
                    col_left = left + c_idx * (col_width + gap)
                    render_elements_to_slide(slide, col['children'], theme, theme_name, col_left, col_width, current_y)
                    
            current_y += col_h + Inches(0.15)
            
        elif elem['type'] == 'table':
            table_data = elem['data']
            if table_data:
                rows = len(table_data)
                cols = len(table_data[0])
                table_shape = slide.shapes.add_table(rows, cols, left, current_y, width, Inches(rows * 0.45))
                table = table_shape.table
                
                for r_idx in range(rows):
                    for c_idx in range(cols):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = table_data[r_idx][c_idx] if c_idx < len(table_data[r_idx]) else ""
                        p = cell.text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.LEFT
                        p.font.name = 'Calibri'
                        p.font.size = Pt(14)
                        p.font.color.rgb = theme['text']
                        
                        if r_idx == 0:
                            p.font.bold = True
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = theme['primary']
                            p.font.color.rgb = RGBColor(255, 255, 255)
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = theme['block_bg']
                            
                current_y += Inches(rows * 0.45) + Inches(0.2)
                
        i += 1

# --- Slide creation wrappers ---
def is_transition_slide(frame):
    return frame['title'].strip().lower() == 'transition'

def create_title_slide(prs, metadata, theme, theme_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme['bg']
    
    current_y = Inches(1.5)
    
    title_text = clean_latex_formatting(metadata.get('title', 'Presentation'))
    title_height = Inches(1.5)
    txBox = slide.shapes.add_textbox(Inches(1.0), current_y, Inches(11.33), title_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Calibri'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme['primary']
    
    current_y += title_height + Inches(0.2)
    
    if 'subtitle' in metadata:
        sub_text = clean_latex_formatting(metadata['subtitle'])
        sub_height = Inches(0.8)
        txBox = slide.shapes.add_textbox(Inches(1.0), current_y, Inches(11.33), sub_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sub_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = theme['text_muted']
        current_y += sub_height + Inches(0.5)
    else:
        current_y += Inches(0.8)
        
    footer_text = ""
    if 'author' in metadata:
        footer_text += clean_latex_formatting(metadata['author'])
    if 'institute' in metadata:
        footer_text += f"\n{clean_latex_formatting(metadata['institute'])}"
    if 'date' in metadata:
        footer_text += f"\n{clean_latex_formatting(metadata['date'])}"
        
    if footer_text:
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.33), Inches(1.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = footer_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(16)
        p.font.color.rgb = theme['text']

def create_transition_slide(prs, frame, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme['block_bg']
    
    block_title = ""
    match = re.search(r'\\begin{block}{([^}]+)}', frame['body'])
    if match:
        block_title = clean_latex_formatting(match.group(1))
        
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = block_title if block_title else "Transition"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Calibri'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = theme['primary']

def create_content_slide(prs, frame, theme, theme_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme['bg']
    
    title_height = Inches(0.8)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), title_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = frame['title']
    p.font.name = 'Calibri'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme['primary']
    
    elements = parse_elements(frame['body'])
    render_elements_to_slide(slide, elements, theme, theme_name, Inches(0.8), Inches(11.73), Inches(1.6))

# --- Main Compiler CLI ---
def main():
    parser = argparse.ArgumentParser(description="Convert LaTeX Beamer to PowerPoint")
    parser.add_argument('tex_file', help="Path to Beamer LaTeX file")
    parser.add_argument('--output', '-o', help="Output PowerPoint path")
    parser.add_argument('--theme', '-t', choices=['dark', 'light'], default='dark', help="Color theme (dark/light)")
    args = parser.parse_args()
    
    if not os.path.exists(args.tex_file):
        print(f"Error: File not found: {args.tex_file}")
        sys.exit(1)
        
    with open(args.tex_file, 'r', encoding='utf-8') as f:
        latex_content = f.read()
        
    # Strip single-line comments (ignoring escaped ones)
    lines = latex_content.splitlines()
    cleaned_lines = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith('%'):
            continue
        # Find unescaped %
        pos = 0
        comment_idx = -1
        while True:
            idx = line.find('%', pos)
            if idx == -1:
                break
            if idx > 0 and line[idx-1] == '\\':
                pos = idx + 1
            else:
                comment_idx = idx
                break
        if comment_idx != -1:
            line = line[:comment_idx]
        cleaned_lines.append(line)
    latex_content = "\n".join(cleaned_lines)
    
    # Parse metadata
    metadata = {}
    for meta_name in ['title', 'subtitle', 'author', 'institute', 'date']:
        # Match e.g., \title[Short]{Long} or \title{Long}
        match = re.search(r'\\' + meta_name + r'(?:\[[^\]]*\])?{', latex_content)
        if match:
            open_brace = match.end() - 1
            content, _ = find_curly_content(latex_content, open_brace)
            if content:
                metadata[meta_name] = content
                
    # Parse frames
    frames = []
    pos = 0
    while True:
        start_match = latex_content.find(r'\begin{frame}', pos)
        if start_match == -1:
            break
        end_match = latex_content.find(r'\end{frame}', start_match)
        if end_match == -1:
            break
            
        header_end = start_match + len(r'\begin{frame}')
        while header_end < len(latex_content) and latex_content[header_end].isspace():
            header_end += 1
        if header_end < len(latex_content) and latex_content[header_end] == '[':
            bracket_count = 1
            header_end += 1
            while header_end < len(latex_content) and bracket_count > 0:
                if latex_content[header_end] == '[':
                    bracket_count += 1
                elif latex_content[header_end] == ']':
                    bracket_count -= 1
                header_end += 1
                
        while header_end < len(latex_content) and latex_content[header_end].isspace():
            header_end += 1
            
        title = ""
        content_start = header_end
        if header_end < len(latex_content) and latex_content[header_end] == '{':
            title_content, next_pos = find_curly_content(latex_content, header_end)
            if title_content is not None:
                title = title_content
                content_start = next_pos
                
        frame_body = latex_content[content_start:end_match].strip()
        frames.append({
            'title': clean_latex_formatting(title),
            'raw_title': title,
            'body': frame_body
        })
        pos = end_match + len(r'\end{frame}')
        
    print(f"Parsed metadata: {list(metadata.keys())}")
    print(f"Parsed {len(frames)} slides.")
    
    # Initialize PowerPoint
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    theme = THEMES[args.theme]
    
    # Create Title Slide
    create_title_slide(prs, metadata, theme, args.theme)
    
    # Create Content and Transition Slides
    for idx, frame in enumerate(frames):
        # Skip the title slide frame itself if it contains \titlepage
        if r'\titlepage' in frame['body']:
            continue
        print(f"Generating slide {idx+1}: {frame['title']}")
        if is_transition_slide(frame):
            create_transition_slide(prs, frame, theme)
        else:
            create_content_slide(prs, frame, theme, args.theme)
            
    output_path = args.output
    if not output_path:
        # Default to same directory with pptx extension
        output_path = os.path.splitext(args.tex_file)[0] + ".pptx"
        
    prs.save(output_path)
    print(f"Successfully saved PowerPoint to: {output_path}")
    print(f"File exists right after save: {os.path.exists(output_path)}")
    if os.path.exists(output_path):
        print(f"File size right after save: {os.path.getsize(output_path)} bytes")

if __name__ == "__main__":
    main()
